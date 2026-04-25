from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
import os

os.chdir(r'c:\Programming\Github\Python')
FIG_DIR = './SWATCH/ppt_figs'

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

TITLE_BG   = RGBColor(0x1F, 0x49, 0x7D)
SECTION_BG = RGBColor(0x2E, 0x75, 0xB6)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK       = RGBColor(0x26, 0x26, 0x26)
ACCENT     = RGBColor(0xED, 0x7D, 0x31)
LIGHT_BG   = RGBColor(0xF2, 0xF7, 0xFF)

blank_layout = prs.slide_layouts[6]

def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.line.width = line_width
    if fill_color:
        shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width if line_width else Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, text, left, top, width, height,
                font_size=14, bold=False, color=DARK, align=PP_ALIGN.LEFT,
                font_name='맑은 고딕', wrap=True):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.size = Pt(font_size); run.font.bold = bold
    run.font.color.rgb = color; run.font.name = font_name
    return txBox

def add_image(slide, path, left, top, width, height=None):
    if height:
        slide.shapes.add_picture(path, Inches(left), Inches(top), Inches(width), Inches(height))
    else:
        slide.shapes.add_picture(path, Inches(left), Inches(top), width=Inches(width))

# ═══════════════════════════════════════════════════════════
# Slide 1: Title
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, fill_color=TITLE_BG)
add_rect(slide, 0, 4.5, 13.33, 0.08, fill_color=ACCENT)
add_rect(slide, 0, 4.58, 13.33, 2.92, fill_color=RGBColor(0x17, 0x37, 0x5E))

add_textbox(slide, 'SWATCH 활성탄소섬유(ACF) 보호능력 분석',
            0.6, 1.0, 12.0, 1.2, font_size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide, '변수 변환 선형회귀 / 다항 Ridge 회귀 / 가우시안 프로세스 회귀',
            0.6, 2.3, 12.0, 0.8, font_size=20, bold=False, color=RGBColor(0xBD, 0xD7, 0xEE), align=PP_ALIGN.CENTER)
add_textbox(slide, '■  GD / HD 보호농도 예측 모형 비교 연구',
            0.6, 3.2, 12.0, 0.6, font_size=16, color=RGBColor(0x9D, 0xC3, 0xE6), align=PP_ALIGN.CENTER)

add_textbox(slide, '분석 대상: ACF25-01 ~ ACF25-09  (n = 9)',
            1.0, 5.0, 6.0, 0.5, font_size=14, color=RGBColor(0x9D, 0xC3, 0xE6))
add_textbox(slide, '예측 변수: 무게(g/m²), BET 비표면적(m²/g)',
            1.0, 5.5, 7.0, 0.5, font_size=14, color=RGBColor(0x9D, 0xC3, 0xE6))
add_textbox(slide, '목표 변수: GD 보호농도(기준 ≤ 357), HD 보호농도(기준 ≤ 671)',
            1.0, 6.0, 10.0, 0.5, font_size=14, color=RGBColor(0x9D, 0xC3, 0xE6))

# ═══════════════════════════════════════════════════════════
# Slide 2: Data Overview
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 0.9, fill_color=SECTION_BG)
add_textbox(slide, '1. 데이터 개요 및 합격/불합격 현황',
            0.3, 0.1, 12.0, 0.7, font_size=22, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_rect(slide, 0, 0.9, 13.33, 6.6, fill_color=LIGHT_BG)

# data table
table_data = [
    ['Sample',    '무게(g/m²)', 'BET(m²/g)', 'GD',  'HD',  'GD 판정', 'HD 판정'],
    ['ACF25-01',  '105.2',     '1150',       '320', '280', '합격',    '합격'],
    ['ACF25-02',  '99.3',      '1450',       '470', '220', '불합격',  '합격'],
    ['ACF25-03',  '119.7',     '1150',       '1150','550', '불합격',  '합격'],
    ['ACF25-04',  '110.4',     '1300',       '850', '850', '불합격',  '불합격'],
    ['ACF25-05',  '103.4',     '1016',       '470', '580', '불합격',  '합격'],
    ['ACF25-06',  '122.8',     '1080',       '20',  '140', '합격',    '합격'],
    ['ACF25-07',  '154.5',     '1139',       '380', '10',  '불합격',  '합격'],
    ['ACF25-08',  '140.0',     '1300',       '20',  '5',   '합격',    '합격'],
    ['ACF25-09',  '130.0',     '2197',       '50',  '5',   '합격',    '합격'],
]
col_w = [1.2, 1.2, 1.1, 0.7, 0.7, 1.0, 1.0]
col_x = [0.15]
for cw in col_w[:-1]: col_x.append(col_x[-1]+cw)
row_h = 0.38; row_y0 = 1.0
for r, row in enumerate(table_data):
    for c, (cell, cx, cw) in enumerate(zip(row, col_x, col_w)):
        is_header = (r == 0)
        is_pass = ('합격' == cell and not is_header)
        is_fail = ('불합격' == cell and not is_header)
        bg = SECTION_BG if is_header else (RGBColor(0xC6,0xEF,0xCE) if is_pass else
                                           (RGBColor(0xFF,0xC7,0xCE) if is_fail else WHITE))
        add_rect(slide, cx, row_y0 + r*row_h, cw-0.02, row_h-0.02,
                 fill_color=bg, line_color=RGBColor(0xCC,0xCC,0xCC), line_width=Pt(0.5))
        fc = WHITE if is_header else DARK
        add_textbox(slide, cell, cx+0.02, row_y0+r*row_h, cw-0.04, row_h-0.02,
                    font_size=10, bold=is_header, color=fc, align=PP_ALIGN.CENTER)

# add figure
add_image(slide, f'{FIG_DIR}/fig1_data.png', 7.2, 1.0, 5.9)
add_textbox(slide, '※ GD 기준 ≤ 357, HD 기준 ≤ 671 (mg·min/m³)',
            0.2, 6.7, 7.0, 0.5, font_size=10, color=RGBColor(0x70,0x70,0x70))

# ═══════════════════════════════════════════════════════════
# Slide 3: Correlation Analysis
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 0.9, fill_color=SECTION_BG)
add_textbox(slide, '2. 상관관계 분석 — 원본 vs 로그 변환',
            0.3, 0.1, 12.0, 0.7, font_size=22, bold=True, color=WHITE)
add_rect(slide, 0, 0.9, 13.33, 6.6, fill_color=LIGHT_BG)

add_image(slide, f'{FIG_DIR}/fig2_corr.png', 0.3, 1.0, 12.5)

add_textbox(slide, '◆ 핵심 관찰',
            0.4, 6.0, 3.0, 0.4, font_size=12, bold=True, color=SECTION_BG)
add_textbox(slide, '로그 변환 후 log(HD)↔log(무게): -0.83,  log(HD)↔log(BET): -0.50  (상관도 크게 향상)',
            0.4, 6.35, 12.5, 0.5, font_size=11, color=DARK)

# ═══════════════════════════════════════════════════════════
# Slide 4: Transform Comparison
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 0.9, fill_color=SECTION_BG)
add_textbox(slide, '3. 변수 변환별 선형회귀 성능 비교',
            0.3, 0.1, 12.0, 0.7, font_size=22, bold=True, color=WHITE)
add_rect(slide, 0, 0.9, 13.33, 6.6, fill_color=LIGHT_BG)

add_image(slide, f'{FIG_DIR}/fig3_transform.png', 0.3, 1.0, 12.5)

add_textbox(slide, '◆ 결론: GD 모형은 어떤 변환에서도 F p > 0.05 (유의하지 않음)',
            0.4, 6.1, 9.0, 0.4, font_size=11, bold=True, color=RGBColor(0xC0,0x00,0x00))
add_textbox(slide, 'HD 모형은 log(HD) ~ 무게 + log(BET) 변환 시 adj-R²=0.775, F p=0.0048로 유의',
            0.4, 6.5, 12.0, 0.5, font_size=11, bold=True, color=RGBColor(0x37,0x5A,0x23))

# ═══════════════════════════════════════════════════════════
# Slide 5: Best HD Model Detail
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 0.9, fill_color=SECTION_BG)
add_textbox(slide, '4. HD 최적 변환 OLS 모형 상세 결과',
            0.3, 0.1, 12.0, 0.7, font_size=22, bold=True, color=WHITE)
add_rect(slide, 0, 0.9, 13.33, 6.6, fill_color=LIGHT_BG)

add_image(slide, f'{FIG_DIR}/fig4_HD_best.png', 0.3, 1.0, 12.5)

# key stats box
add_rect(slide, 0.3, 5.6, 12.5, 1.6, fill_color=WHITE,
         line_color=SECTION_BG, line_width=Pt(1.5))
stats_text = [
    ('모형: log(HD) ~ 무게 + log(BET)',                     0.5, 5.65, 6.0),
    ('R² = 0.831,  adj-R² = 0.775',                         0.5, 6.05, 6.0),
    ('F 통계량 p-값 = 0.0048 (유의수준 1% 이내)',           6.5, 5.65, 6.0),
    ('p(무게) = 0.004  ★★   p(BET) = 0.043  ★',           6.5, 6.05, 6.0),
]
for text, lx, ty, tw in stats_text:
    add_textbox(slide, text, lx, ty, tw, 0.45, font_size=12, bold=True, color=SECTION_BG)

# ═══════════════════════════════════════════════════════════
# Slide 6: Ridge Regression
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 0.9, fill_color=SECTION_BG)
add_textbox(slide, '5. 다항 Ridge 회귀 (2차항 + 상호작용항)',
            0.3, 0.1, 12.0, 0.7, font_size=22, bold=True, color=WHITE)
add_rect(slide, 0, 0.9, 13.33, 6.6, fill_color=LIGHT_BG)

add_image(slide, f'{FIG_DIR}/fig5_ridge.png', 0.3, 1.0, 12.5)

add_textbox(slide,
    '특성: 무게, BET, 무게², BET², 무게×BET  →  StandardScaler 표준화  →  RidgeCV (LOOCV 기반 α 선택)',
    0.4, 5.9, 12.5, 0.45, font_size=11, color=DARK)
add_textbox(slide,
    '※ n=9, 특성 5개: LOO 시 훈련 데이터 8개로 5개 특성 적합 → 과적합 위험. 학습 R² > LOO R²',
    0.4, 6.35, 12.5, 0.5, font_size=11, color=RGBColor(0xC0,0x00,0x00))

# ═══════════════════════════════════════════════════════════
# Slide 7: GPR
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 0.9, fill_color=SECTION_BG)
add_textbox(slide, '6. 가우시안 프로세스 회귀 (GPR)',
            0.3, 0.1, 12.0, 0.7, font_size=22, bold=True, color=WHITE)
add_rect(slide, 0, 0.9, 13.33, 6.6, fill_color=LIGHT_BG)

add_image(slide, f'{FIG_DIR}/fig6_gpr.png', 0.3, 1.0, 12.5)

add_textbox(slide,
    '커널: C(1.0) × RBF([l₁, l₂]) + WhiteKernel(σ²)   |   n_restarts=15로 커널 파라미터 최적화   |   normalize_y=True',
    0.4, 5.9, 12.5, 0.45, font_size=11, color=DARK)
add_textbox(slide,
    '※ GD: 학습 R²=1.000 (보간 과적합), LOO R²<0  /  HD: 학습 R²=0.68, LOO R²<0  →  소표본 한계 명확',
    0.4, 6.35, 12.5, 0.5, font_size=11, color=RGBColor(0xC0,0x00,0x00))

# ═══════════════════════════════════════════════════════════
# Slide 8: Model Summary Comparison
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 0.9, fill_color=SECTION_BG)
add_textbox(slide, '7. 모형 성능 종합 비교 (LOO Cross-Validation)',
            0.3, 0.1, 12.0, 0.7, font_size=22, bold=True, color=WHITE)
add_rect(slide, 0, 0.9, 13.33, 6.6, fill_color=LIGHT_BG)

add_image(slide, f'{FIG_DIR}/fig7_summary.png', 0.3, 1.0, 12.5)

# summary table
tbl_data = [
    ['모형',          'GD 학습 R²', 'GD LOO R²', 'GD LOO RMSE', 'HD 학습 R²', 'HD LOO R²', 'HD LOO RMSE'],
    ['원본 OLS',      '0.182',     '-',          '-',           '0.454',      '-',          '-'],
    ['변환 OLS (최적)','0.243(GD)',  '-',          '-',           '0.831★★',   '-',          '-'],
    ['다항 Ridge',    '0.175',     '-0.235',     '~500',        '0.419',      '0.134',      '~280'],
    ['GPR',           '1.000',     '-0.266',     '~409',        '0.680',      '-0.199',     '~313'],
]
cw2 = [2.0, 1.4, 1.4, 1.5, 1.4, 1.4, 1.5]
cx2 = [0.2]
for c in cw2[:-1]: cx2.append(cx2[-1]+c)
ry0 = 5.5; rh = 0.36
for r, row in enumerate(tbl_data):
    for c, (cell, lx, cw) in enumerate(zip(row, cx2, cw2)):
        is_hdr = (r == 0)
        is_best = ('★' in cell)
        bg = SECTION_BG if is_hdr else (RGBColor(0xC6,0xEF,0xCE) if is_best else WHITE)
        add_rect(slide, lx, ry0+r*rh, cw-0.02, rh-0.02,
                 fill_color=bg, line_color=RGBColor(0xCC,0xCC,0xCC), line_width=Pt(0.5))
        fc = WHITE if is_hdr else (RGBColor(0x37,0x5A,0x23) if is_best else DARK)
        add_textbox(slide, cell, lx+0.02, ry0+r*rh, cw-0.04, rh-0.02,
                    font_size=9, bold=(is_hdr or is_best), color=fc, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════
# Slide 9: Conclusion
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 0.9, fill_color=TITLE_BG)
add_textbox(slide, '8. 결론 및 시사점',
            0.3, 0.1, 12.0, 0.7, font_size=24, bold=True, color=WHITE)
add_rect(slide, 0, 0.9, 13.33, 6.6, fill_color=LIGHT_BG)

conclusions = [
    ('GD 보호농도',
     '무게(g/m²)와 BET 비표면적만으로는 GD 농도를 유의하게 설명하지 못함 (모든 모형 F p > 0.05).\n추가 설명 변수(원단 구조, 처리방법 등) 필요.',
     RGBColor(0xC0,0x00,0x00)),
    ('HD 보호농도 ★ 핵심 발견',
     'log(HD) ~ 무게 + log(BET) 변환 OLS가 R²=0.831, adj-R²=0.775, F p=0.0048로 고도로 유의.\n→ HD 보호능력은 무게↑ + BET↑ 일수록 낮아지는(보호능력↑) 강한 선형 관계 존재.',
     RGBColor(0x37,0x5A,0x23)),
    ('다항 Ridge 회귀',
     'n=9 대비 5개 특성으로 LOO 교차검증 시 음의 R² 발생 (과적합).\n소표본에서는 단순 선형모형이 더 안정적. 변환 OLS HD 모형 권장.',
     SECTION_BG),
    ('가우시안 프로세스 회귀',
     'GD 학습 R²=1.0 (보간) 이지만 LOO R²<0 — 외삽 예측 불안정.\n소표본 비선형 탐색 목적으로는 유용하나 예측 모형으로는 부적합.',
     SECTION_BG),
    ('향후 방향',
     '추가 샘플 확보 (최소 n≥20 권장) → Ridge/GPR 예측력 재평가.\nGD 영향 인자 추가 실험 설계 필요.',
     RGBColor(0x70,0x30,0xA0)),
]

for i, (title, body, color) in enumerate(conclusions):
    ty = 1.05 + i * 1.12
    add_rect(slide, 0.3, ty, 0.1, 0.85, fill_color=color)
    add_textbox(slide, title, 0.5, ty, 3.5, 0.45, font_size=13, bold=True, color=color)
    add_textbox(slide, body, 0.5, ty+0.38, 12.3, 0.72, font_size=11, color=DARK, wrap=True)

# ─── Save ───────────────────────────────────────────────
out_path = './SWATCH/SWATCH_분석결과_발표자료.pptx'
prs.save(out_path)
print(f'PPT saved: {out_path}')
